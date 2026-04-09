#!/usr/bin/env /opt/homebrew/bin/python3.11
"""
Generate a polished PowerPoint (.pptx) from a structured YAML slide spec.

Usage:
    python3.11 src/utils/generate_pptx.py <input_yaml> <output_pptx>
"""

import sys
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


def hex_to_rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_speaker_notes(slide, notes_text, font_name):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text
    for p in tf.paragraphs:
        p.font.name = font_name
        p.font.size = Pt(11)


def add_accent_bar(slide, left, top, width, height, color):
    """Add a visible accent bar with rounded corners."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Reduce corner rounding
    shape.adjustments[0] = 0.1
    return shape


def add_side_stripe(slide, color, width=Inches(0.35)):
    """Add a vertical accent stripe on the left edge."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), width, Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def make_para(tf, text, font_name, font_size, color, bold=False, italic=False,
              alignment=PP_ALIGN.LEFT, space_before=0, space_after=0, indent=False):
    """Add a paragraph to a text frame with full formatting."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if indent:
        p.level = 1
    return p


def add_body_text(tf, lines, font_name, font_size, color, sub_color=None, accent_color=None):
    """Add body lines with proper hierarchy.

    Line conventions:
      **text**  → bold emphasis line (larger spacing above)
      - text    → indented sub-bullet (italic, smaller)
      text      → normal body line
    """
    tf.word_wrap = True
    if sub_color is None:
        sub_color = color
    for i, raw_line in enumerate(lines):
        line = raw_line.strip()
        if not line:
            continue

        # Bold emphasis lines: **text**
        if line.startswith("**") and line.endswith("**"):
            line = line[2:-2]
            make_para(
                tf, line,
                font_name, font_size,
                color,
                bold=True,
                space_before=16,
                space_after=4,
            )
        # Sub-bullets: - text
        elif line.startswith("- "):
            line = line[2:]
            make_para(
                tf, "     " + line,
                font_name, font_size - 2,
                sub_color,
                space_before=2,
                space_after=6,
                italic=True,
            )
        # Normal body lines
        else:
            make_para(
                tf, line,
                font_name, font_size,
                color,
                space_before=8,
                space_after=4,
            )


# ─── SLIDE BUILDERS ────────────────────────────────────────

def make_title_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_heading"]))

    # Large accent bar at top
    add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(3.0), Pt(6),
                   hex_to_rgb(theme["color_accent"]))

    # Title — large, white, bold
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    make_para(tf, spec.get("title", ""), theme["font_heading"], 44,
              RGBColor(0xFF, 0xFF, 0xFF), bold=True, space_after=8)

    # Subtitle — lighter, smaller
    if spec.get("subtitle"):
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(8.4), Inches(1.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        make_para(tf2, spec["subtitle"], theme["font_body"], 20,
                  RGBColor(0xCA, 0xDC, 0xFC), space_after=4)

    # Thin bottom line
    add_accent_bar(slide, Inches(0.8), Inches(6.6), Inches(8.4), Pt(2),
                   hex_to_rgb(theme["color_accent"]))

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


def make_section_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_heading"]))

    # Accent bar
    add_accent_bar(slide, Inches(0.8), Inches(2.8), Inches(2.5), Pt(5),
                   hex_to_rgb(theme["color_accent"]))

    # Section title — centered vertically, large
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(8.4), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    make_para(tf, spec.get("title", ""), theme["font_heading"], 38,
              RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


def make_content_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_bg"]))

    # Left accent stripe
    add_side_stripe(slide, hex_to_rgb(theme["color_heading"]))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.8), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    make_para(tf, spec.get("title", ""), theme["font_heading"], 32,
              hex_to_rgb(theme["color_heading"]), bold=True)

    # Accent bar under title
    add_accent_bar(slide, Inches(0.8), Inches(1.3), Inches(2.0), Pt(4),
                   hex_to_rgb(theme["color_accent"]))

    # Body
    body_text = spec.get("body", "")
    lines = [l for l in body_text.strip().split("\n") if l.strip()]
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.6), Inches(5.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    add_body_text(tf2, lines, theme["font_body"], 20,
                  hex_to_rgb(theme["color_body"]),
                  hex_to_rgb(theme.get("color_sub", theme["color_body"])))

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


def make_two_column_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_bg"]))

    # Left accent stripe
    add_side_stripe(slide, hex_to_rgb(theme["color_heading"]))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.8), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    make_para(tf, spec.get("title", ""), theme["font_heading"], 30,
              hex_to_rgb(theme["color_heading"]), bold=True)

    # Accent bar under title
    add_accent_bar(slide, Inches(0.8), Inches(1.3), Inches(2.0), Pt(4),
                   hex_to_rgb(theme["color_accent"]))

    # Vertical divider line between columns
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.85), Inches(1.6), Pt(1.5), Inches(5.2)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    divider.line.fill.background()

    # Left column
    left_lines = [l for l in spec.get("left", "").strip().split("\n") if l.strip()]
    txBox_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(3.9), Inches(5.3))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    add_body_text(tf_l, left_lines, theme["font_body"], 18,
                  hex_to_rgb(theme["color_body"]))

    # Right column
    right_lines = [l for l in spec.get("right", "").strip().split("\n") if l.strip()]
    txBox_r = slide.shapes.add_textbox(Inches(5.1), Inches(1.65), Inches(4.5), Inches(5.3))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    add_body_text(tf_r, right_lines, theme["font_body"], 18,
                  hex_to_rgb(theme["color_body"]))

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


def make_table_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_bg"]))
    add_side_stripe(slide, hex_to_rgb(theme["color_heading"]))

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.8), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    make_para(tf, spec.get("title", ""), theme["font_heading"], 32,
              hex_to_rgb(theme["color_heading"]), bold=True)

    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    n_rows = len(rows) + 1
    n_cols = len(headers)

    row_h = Inches(0.55)
    table_h = row_h * n_rows
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.8), Inches(1.7), Inches(8.6), table_h
    )
    table = table_shape.table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = hex_to_rgb(theme["color_heading"])
        for p in cell.text_frame.paragraphs:
            p.font.name = theme["font_heading"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.name = theme["font_body"]
                p.font.size = Pt(15)
                p.font.color.rgb = hex_to_rgb(theme["color_body"])

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


def make_quote_slide(prs, spec, theme):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, hex_to_rgb(theme["color_bg"]))
    add_side_stripe(slide, hex_to_rgb(theme["color_heading"]))

    # Title
    if spec.get("title"):
        txBox_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.35), Inches(8.8), Inches(0.9))
        tf_t = txBox_t.text_frame
        tf_t.word_wrap = True
        make_para(tf_t, spec["title"], theme["font_heading"], 32,
                  hex_to_rgb(theme["color_heading"]), bold=True)

    # Large open-quote mark
    txBox_q_mark = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(1.0), Inches(1.0))
    tf_qm = txBox_q_mark.text_frame
    make_para(tf_qm, "\u201C", theme["font_heading"], 72,
              hex_to_rgb(theme["color_accent"]), bold=True)

    # Quote text
    quote_top = Inches(2.4)
    txBox_q = slide.shapes.add_textbox(Inches(1.5), quote_top, Inches(7.5), Inches(2.5))
    tf_q = txBox_q.text_frame
    tf_q.word_wrap = True
    make_para(tf_q, spec.get("quote", ""), theme["font_heading"], 26,
              hex_to_rgb(theme["color_heading"]), italic=True,
              alignment=PP_ALIGN.LEFT, space_after=12)

    # Attribution
    if spec.get("attribution"):
        txBox_a = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(7.5), Inches(0.8))
        tf_a = txBox_a.text_frame
        tf_a.word_wrap = True
        make_para(tf_a, spec["attribution"], theme["font_body"], 16,
                  hex_to_rgb(theme["color_accent"]), italic=True)

    add_speaker_notes(slide, spec.get("notes", ""), theme["font_body"])


SLIDE_BUILDERS = {
    "title": make_title_slide,
    "section": make_section_slide,
    "content": make_content_slide,
    "two_column": make_two_column_slide,
    "table": make_table_slide,
    "quote": make_quote_slide,
}


def generate(input_path, output_path):
    with open(input_path, "r") as f:
        spec = yaml.safe_load(f)

    theme = spec.get("theme", {})
    theme.setdefault("font_heading", "Georgia")
    theme.setdefault("font_body", "Calibri")
    theme.setdefault("color_heading", "1E2761")
    theme.setdefault("color_accent", "7B2D26")
    theme.setdefault("color_body", "2c3e50")
    theme.setdefault("color_bg", "f7f6f2")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_spec in spec.get("slides", []):
        slide_type = slide_spec.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type, make_content_slide)
        builder(prs, slide_spec, theme)

    prs.save(output_path)
    print(f"Generated: {output_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: generate_pptx.py <input.yaml> <output.pptx>")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2])
